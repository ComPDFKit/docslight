"""Office document loading as Markdown."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from docslight.exceptions import DependencyMissingError, UnsupportedFormatError
from docslight.local.loaders import LOCAL_DEPS_MESSAGE, LoadedTextDocument

LEGACY_OFFICE_EXTENSIONS = {".doc", ".ppt", ".xls"}


def load_workbook(path: Path, **kwargs: Any) -> Any:
    """Load an XLSX workbook with optional dependency handling."""
    try:
        from openpyxl import load_workbook as openpyxl_load_workbook
    except ModuleNotFoundError as exc:  # pragma: no cover - depends on environment
        raise DependencyMissingError(LOCAL_DEPS_MESSAGE) from exc
    return openpyxl_load_workbook(path, **kwargs)


class OfficeMarkdownLoader:
    """Load modern Office files into basic Markdown."""

    def load(self, path: Path | str) -> LoadedTextDocument:
        """Load a DOCX, PPTX, or XLSX file into Markdown."""
        source_path = Path(path)
        suffix = source_path.suffix.lower()
        if suffix in LEGACY_OFFICE_EXTENSIONS:
            raise UnsupportedFormatError("Legacy Office files must convert to DOCX, PPTX, or XLSX")
        if suffix == ".docx":
            return self._load_docx(source_path)
        if suffix == ".pptx":
            return self._load_pptx(source_path)
        if suffix == ".xlsx":
            return self._load_xlsx(source_path)
        raise UnsupportedFormatError(f"Unsupported Office format: {suffix or source_path.name}")

    def _load_docx(self, path: Path) -> LoadedTextDocument:
        try:
            from docx import Document
        except ModuleNotFoundError as exc:  # pragma: no cover - depends on environment
            raise DependencyMissingError(LOCAL_DEPS_MESSAGE) from exc

        document = Document(str(path))
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            rows: list[list[str]] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(cells)
            if rows:
                parts.append(self._rows_to_markdown(rows))
        return LoadedTextDocument(
            markdown="\n\n".join(parts),
            metadata={
                "document_type": "docx",
                "page_count": 1,
                "paragraph_count": len(document.paragraphs),
                "table_count": len(document.tables),
            },
        )

    def _load_pptx(self, path: Path) -> LoadedTextDocument:
        try:
            from pptx import Presentation
        except ModuleNotFoundError as exc:  # pragma: no cover - depends on environment
            raise DependencyMissingError(LOCAL_DEPS_MESSAGE) from exc

        presentation = Presentation(str(path))
        slide_parts: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text.strip():
                    texts.append(text.strip())
            slide_parts.append("\n\n".join([f"# Slide {index}", *texts]))
        return LoadedTextDocument(
            markdown="\n\n".join(slide_parts),
            metadata={
                "document_type": "pptx",
                "page_count": len(presentation.slides),
                "slide_count": len(presentation.slides),
            },
        )

    def _load_xlsx(self, path: Path) -> LoadedTextDocument:
        workbook = load_workbook(path, data_only=True, read_only=True)
        sheets: list[str] = []
        try:
            for worksheet in workbook.worksheets:
                rows = [self._format_row(row) for row in worksheet.iter_rows(values_only=True)]
                rows = [row for row in rows if row]
                sheets.append(
                    "\n".join([f"## Sheet: {worksheet.title}", self._rows_to_markdown(rows)])
                )
            return LoadedTextDocument(
                markdown="\n\n".join(sheets),
                metadata={
                    "document_type": "xlsx",
                    "page_count": len(workbook.worksheets),
                    "sheet_count": len(workbook.worksheets),
                    "sheet_names": [worksheet.title for worksheet in workbook.worksheets],
                },
            )
        finally:
            workbook.close()

    def _format_row(self, row: tuple[Any, ...]) -> list[str]:
        values = ["" if value is None else str(value) for value in row]
        if not any(values):
            return []
        return values

    def _rows_to_markdown(self, rows: list[list[str]]) -> str:
        if not rows:
            return ""
        escaped_rows = [[self._escape_table_cell(cell) for cell in row] for row in rows]
        header = escaped_rows[0]
        separator = ["---"] * len(header)
        markdown_rows = [header, separator, *escaped_rows[1:]]
        return "\n".join("| " + " | ".join(row) + " |" for row in markdown_rows)

    def _escape_table_cell(self, value: str) -> str:
        return value.replace("\n", " ").replace("\r", " ").replace("|", "\\|")
