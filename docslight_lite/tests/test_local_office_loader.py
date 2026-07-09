from __future__ import annotations

from pathlib import Path

import pytest

from docslight.exceptions import UnsupportedFormatError
from docslight.local import OfficeMarkdownLoader


def test_office_loader_loads_real_docx_with_text_table_and_metadata(
    tmp_path: Path,
) -> None:
    from docx import Document

    path = tmp_path / "invoice.docx"
    document = Document()
    document.add_heading("Invoice Summary", level=1)
    document.add_paragraph("Payment is due on receipt.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Invoice"
    table.cell(0, 1).text = "Total"
    table.cell(1, 0).text = "INV-001"
    table.cell(1, 1).text = "100"
    document.save(path)

    result = OfficeMarkdownLoader().load(path)

    assert "Invoice Summary" in result.markdown
    assert "Payment is due on receipt." in result.markdown
    assert "| Invoice | Total |" in result.markdown
    assert "| --- | --- |" in result.markdown
    assert result.metadata["document_type"] == "docx"
    assert result.metadata["table_count"] == 1
    assert result.metadata["paragraph_count"] >= 2


def test_office_loader_loads_real_pptx_with_slide_text_and_metadata(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Quarterly Invoice Review"
    presentation.save(path)

    result = OfficeMarkdownLoader().load(path)

    assert "# Slide 1" in result.markdown
    assert "Quarterly Invoice Review" in result.markdown
    assert result.metadata["document_type"] == "pptx"
    assert result.metadata["slide_count"] == 1
    assert result.metadata["page_count"] == 1


def test_office_loader_loads_real_xlsx_with_table_and_metadata(tmp_path: Path) -> None:
    from openpyxl import Workbook

    path = tmp_path / "invoices.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Invoices"
    worksheet.append(["Invoice", "Total"])
    worksheet.append(["INV-001", 100])
    workbook.save(path)

    result = OfficeMarkdownLoader().load(path)

    assert "## Sheet: Invoices" in result.markdown
    assert "| Invoice | Total |" in result.markdown
    assert "| --- | --- |" in result.markdown
    assert "| INV-001 | 100 |" in result.markdown
    assert result.metadata["document_type"] == "xlsx"
    assert result.metadata["sheet_count"] == 1
    assert result.metadata["sheet_names"] == ["Invoices"]
    assert result.metadata["page_count"] == 1


def test_office_loader_escapes_xlsx_table_cell_pipes_and_newlines(
    tmp_path: Path,
) -> None:
    from openpyxl import Workbook

    path = tmp_path / "escaped.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Invoices"
    worksheet.append(["Invoice", "Notes"])
    worksheet.append(["A|B", "line1\nline2"])
    workbook.save(path)

    result = OfficeMarkdownLoader().load(path)

    assert "| A\\|B | line1 line2 |" in result.markdown
    assert "| A|B |" not in result.markdown
    assert "line1\nline2" not in result.markdown


def test_office_loader_rejects_legacy_xls_with_conversion_message(
    tmp_path: Path,
) -> None:
    path = tmp_path / "legacy.xls"
    path.write_bytes(b"legacy")

    with pytest.raises(UnsupportedFormatError, match="convert to DOCX, PPTX, or XLSX"):
        OfficeMarkdownLoader().load(path)
